"""Generate deck/presentation.pptx — the manager-facing MVP deck.

Run any time: `python deck/build_deck.py`. The demo slide pulls the newest
annotated snapshots and the tail of logs/events.csv when they exist, and
renders labeled placeholders when they don't — so the deck always builds.
Regenerate after the staged demo run to pull in real screenshots.

Slide content spec: docs/mvp-plan.md §5.
"""

from __future__ import annotations

import csv
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
LOG_DIR = ROOT / "logs"
SNAP_DIR = LOG_DIR / "snapshots"
EVENTS_CSV = LOG_DIR / "events.csv"
OUT_PATH = ROOT / "deck" / "presentation.pptx"

DARK = RGBColor(0x1F, 0x2A, 0x36)
ACCENT = RGBColor(0x2A, 0x6F, 0x97)
BOX_FILL = RGBColor(0xE8, 0xF0, 0xF7)
WARN_FILL = RGBColor(0xFB, 0xEE, 0xDD)
GREY = RGBColor(0x6B, 0x77, 0x85)


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_title(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.1), Inches(0.9))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK


def add_bullets(slide, items, top=1.4, left=0.8, width=11.7, size=18):
    """items: list of str or (str, level) tuples."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.6))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        text, level = item if isinstance(item, tuple) else (item, 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(size if level == 0 else size - 3)
        p.font.color.rgb = DARK if level == 0 else GREY
        p.space_after = Pt(10)
    return tb


def add_box(slide, text, x, y, w, h, diamond=False, fill=BOX_FILL):
    shape_type = MSO_SHAPE.DIAMOND if diamond else MSO_SHAPE.ROUNDED_RECTANGLE
    box = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = ACCENT
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = DARK
    return box


def _arrowhead(conn):
    ln = conn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))


def connect(slide, a, b, from_idx=3, to_idx=1, elbow=False):
    """Connector from shape a to shape b. Sites: 0=top 1=left 2=bottom 3=right."""
    kind = MSO_CONNECTOR.ELBOW if elbow else MSO_CONNECTOR.STRAIGHT
    conn = slide.shapes.add_connector(kind, 0, 0, 0, 0)
    conn.begin_connect(a, from_idx)
    conn.end_connect(b, to_idx)
    conn.line.color.rgb = GREY
    conn.line.width = Pt(2)
    _arrowhead(conn)
    return conn


def add_label(slide, text, x, y, w=1.8, size=11, italic=True, align=PP_ALIGN.CENTER):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.italic = italic
    p.font.color.rgb = GREY
    return tb


def slide_1_title(prs):
    s = blank_slide(prs)
    tb = s.shapes.add_textbox(Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.4))
    p = tb.text_frame.paragraphs[0]
    p.text = "Storage Unit Item Detection — MVP"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK
    sub = s.shapes.add_textbox(Inches(0.9), Inches(4.0), Inches(11.5), Inches(1.2))
    tf = sub.text_frame
    tf.paragraphs[0].text = "Detect items in view · flag when something is added or removed"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.color.rgb = ACCENT
    p2 = tf.add_paragraph()
    p2.text = f"Misael · {date.today().isoformat()} · one-day MVP on stock models"
    p2.font.size = Pt(16)
    p2.font.color.rgb = GREY


def slide_2_problem(prs):
    s = blank_slide(prs)
    add_title(s, "The problem")
    add_bullets(s, [
        "A storage unit holds changing inventory — nobody watches it continuously",
        "Goal 1: know what items are currently in view",
        "Goal 2: raise an event the moment an item is added or removed",
        "Constraint: commodity hardware — one webcam, CPU-only machine (no GPU)",
        "Deliverable today: the smallest version that demonstrably works, plus the roadmap to a reliable one",
    ])


def slide_3_architecture(prs):
    s = blank_slide(prs)
    add_title(s, "MVP architecture")

    # Row 1: capture → sample → detect → count
    y1 = 1.6
    cam = add_box(s, "Webcam\n(continuous connection)", 0.4, y1, 2.2, 1.1)
    loop = add_box(s, "Frame loop\n(read every frame)", 3.1, y1, 2.2, 1.1)
    sample = add_box(s, "Sampled frame\n(every 5 s)", 5.8, y1, 2.2, 1.1)
    yolo = add_box(s, "YOLO11n on CPU\n(stock COCO classes)", 8.5, y1, 2.2, 1.1)
    counts = add_box(s, "Per-class counts", 11.0, y1, 1.9, 1.1)
    connect(s, cam, loop)
    connect(s, loop, sample)
    connect(s, sample, yolo)
    connect(s, yolo, counts)

    # Row 2: diff → events → output, flowing right-to-left
    y2 = 4.3
    diff = add_box(s, "Diff vs previous\ncounts", 10.3, y2, 2.6, 1.5, diamond=True)
    events = add_box(s, "ADDED / REMOVED\nevents", 6.2, y2 + 0.2, 2.4, 1.1, fill=WARN_FILL)
    out = add_box(s, "Console + events.csv\n+ annotated snapshot", 2.6, y2 + 0.2, 2.6, 1.1)
    connect(s, counts, diff, from_idx=2, to_idx=0, elbow=True)
    connect(s, diff, events, from_idx=1, to_idx=3)
    add_label(s, "change", 9.0, y2 + 0.35)
    connect(s, events, out, from_idx=1, to_idx=3)
    connect(s, out, loop, from_idx=0, to_idx=2, elbow=True)
    add_label(s, "next sample", 2.9, 3.55)
    connect(s, diff, loop, from_idx=0, to_idx=2, elbow=True)
    add_label(s, "no change", 7.0, 3.1)

    add_label(s, "Sampling every N seconds keeps a CPU-only machine comfortably ahead of the detector's speed",
              0.6, 6.7, w=12.0, size=13, align=PP_ALIGN.LEFT)


def slide_4_design_choices(prs):
    s = blank_slide(prs)
    add_title(s, "Key design choices")
    add_bullets(s, [
        "Sampled frames, not full-rate video",
        ("No GPU: YOLO on CPU manages a few frames/second at best — sampling every 5 s from a "
         "continuous connection stays real-time-enough without lag or dropped frames", 1),
        "Stock pretrained YOLO11n — zero training time",
        ("COCO's 80 classes have no \"box\" class: suitcase / backpack / handbag act as stand-ins "
         "for real storage items until Part B custom training", 1),
        "Count-diff event logic — no tracking",
        ("Per class: count went up → ADDED, count went down → REMOVED. Simple, explainable, "
         "unit-testable", 1),
        "Everything logged",
        ("Counts every sample; events to CSV; annotated snapshot saved on every change", 1),
    ])


def _newest_snapshots(n=2):
    if not SNAP_DIR.is_dir():
        return []
    snaps = sorted(SNAP_DIR.glob("*.jpg"), key=lambda p: p.stat().st_mtime, reverse=True)
    return snaps[:n]


def _events_tail(n=8):
    if not EVENTS_CSV.is_file():
        return None
    with EVENTS_CSV.open(newline="") as f:
        rows = list(csv.reader(f))
    if len(rows) <= 1:
        return None
    header, body = rows[0], rows[1:][-n:]
    widths = [max(len(r[i]) for r in [header] + body) for i in range(len(header))]
    fmt = lambda r: "  ".join(c.ljust(w) for c, w in zip(r, widths))
    return "\n".join([fmt(header)] + [fmt(r) for r in body])


def slide_5_demo(prs):
    s = blank_slide(prs)
    add_title(s, "Demo — it works")

    snaps = _newest_snapshots()
    if snaps:
        x = 0.6
        for p in snaps:
            pic = s.shapes.add_picture(str(p), Inches(x), Inches(1.4), height=Inches(3.1))
            add_label(s, p.name, x, 4.55, w=pic.width.inches, size=10)
            x += pic.width.inches + 0.4
    else:
        for i, x in enumerate((0.6, 5.4)):
            ph = add_box(s, "Annotated snapshot placeholder\n\nrun watcher.py, stage an add/remove,\nthen re-run build_deck.py",
                         x, 1.4, 4.4, 3.1)
            ph.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
            ph.line.color.rgb = GREY

    tail = _events_tail()
    tb = s.shapes.add_textbox(Inches(0.6), Inches(5.0), Inches(12.1), Inches(2.2))
    tf = tb.text_frame
    tf.word_wrap = False
    lines = (tail or "events.csv placeholder — rows will appear here after the staged demo run").splitlines()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.name = "Consolas"
        p.font.color.rgb = DARK if tail else GREY


def slide_6_limitations(prs):
    s = blank_slide(prs)
    add_title(s, "Known limitations — stated honestly")
    add_bullets(s, [
        "Stand-in classes: COCO has no \"box\" class — real storage items are approximated until custom training",
        "A single flickered detection produces a false ADDED/REMOVED pair — hysteresis is a Part B item",
        "Equal-count swaps between samples are invisible to count-diff logic",
        "Lighting sensitivity untested — Part B has a dedicated stress-test phase",
        "No tracking across frames: identity of individual items is not maintained",
    ])


def slide_7_roadmap(prs):
    s = blank_slide(prs)
    add_title(s, "Roadmap — Part B (post-MVP)")
    add_bullets(s, [
        "1 · Capture strategy — motion-triggered snapshots (frame diff, MOG2/KNN) vs. intervals; lighting false-trigger testing",
        "2 · Data & annotation — labeling workflow (Roboflow / CVAT), augmentation to stretch a small real dataset",
        "3 · Detection model — fine-tune YOLO on real items, compare one alternative architecture, benchmark on target hardware",
        "4 · Event logic — confidence hysteresis to kill flicker false-events; tracker comparison if continuous video returns",
        "5 · Robustness — lighting, occlusion, camera angle, similar-item confusion, long-running drift",
        "6 · Recommendation — chosen approach, reliability limits, deployment hardware requirements",
    ], size=16)


def slide_8_next_steps(prs):
    s = blank_slide(prs)
    add_title(s, "Next steps / ask")
    add_bullets(s, [
        "Validate the MVP against the real storage unit and real items this week",
        "Decide camera placement and target hardware for deployment",
        "Approve Part B exploration time — capture strategy and custom training are the two biggest reliability levers",
        "Schedule a data-collection session (photos of real items) to unlock custom training",
    ])


def main() -> int:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_1_title(prs)
    slide_2_problem(prs)
    slide_3_architecture(prs)
    slide_4_design_choices(prs)
    slide_5_demo(prs)
    slide_6_limitations(prs)
    slide_7_roadmap(prs)
    slide_8_next_steps(prs)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    n_snaps = len(_newest_snapshots())
    print(f"Wrote {OUT_PATH} ({len(prs.slides)} slides, "
          f"{n_snaps} real snapshot(s), events tail: {'yes' if _events_tail() else 'placeholder'})")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
